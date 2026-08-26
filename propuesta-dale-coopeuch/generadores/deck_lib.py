# -*- coding: utf-8 -*-
"""Sistema de diseño Intothecom para python-pptx. Replica el deck base."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Paleta extraída del PDF base ──────────────────────────────────────────────
ORANGE   = RGBColor(0xE5, 0x70, 0x00)
ORANGE_2 = RGBColor(0xE6, 0x91, 0x38)
BLACK    = RGBColor(0x00, 0x00, 0x00)
INK      = RGBColor(0x26, 0x26, 0x26)
DARK     = RGBColor(0x09, 0x09, 0x09)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_D   = RGBColor(0x59, 0x59, 0x59)
LINE     = RGBColor(0xEE, 0xEE, 0xEE)

FONT = "Be Vietnam Pro"
W, H = Inches(10), Inches(5.625)

ASSETS = "/tmp/claude-0/-home-user-meta-mcp/3417b367-2cff-5bc9-8662-cddb20c827cd/scratchpad"
LOGO_W = f"{ASSETS}/assets/logo_intothecom_blanco.png"
LOGO_B = f"{ASSETS}/assets/logo_intothecom_negro.png"
PHOTO  = f"{ASSETS}/photos"

NAV = ["Paid Media", "Community M.", "Email M.", "Datapify",
       "Casos Éxito", "Propuesta Serv.", "Propuesta Com."]


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(sl, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=Pt(1)):
    s = sl.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = lw
    s.shadow.inherit = False
    return s


def txt(sl, x, y, w, h, runs, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.15, space_after=0, italic=False):
    """runs: str o lista de (texto, {overrides})."""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        parts, seq = runs.split("\n"), []
        for i, part in enumerate(parts):
            if i:
                seq.append(("\n", {}))
            if part:
                seq.append((part, {}))
        runs = seq or [("", {})]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    p.space_after = Pt(space_after)
    for t, ov in runs:
        if t == "\n":
            p = tf.add_paragraph(); p.alignment = align
            p.line_spacing = spacing; p.space_after = Pt(space_after)
            continue
        r = p.add_run(); r.text = t
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
    return tb


def bullets(sl, x, y, w, h, items, size=11, color=GRAY, bullet_color=INK,
            spacing=1.25, gap=7, bold_first=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(gap)
        rb = p.add_run(); rb.text = "•  "
        rb.font.name = FONT; rb.font.size = Pt(size); rb.font.color.rgb = bullet_color
        rb.font.bold = True
        if isinstance(it, tuple):
            head, rest = it
            r1 = p.add_run(); r1.text = head
            r1.font.name = FONT; r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = INK
            r2 = p.add_run(); r2.text = rest
            r2.font.name = FONT; r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = it
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold_first
    return tb


def picture_cover(sl, path, x, y, w, h):
    """Inserta imagen recortada tipo 'cover' dentro del rect (x,y,w,h)."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    pic = sl.shapes.add_picture(path, x, y, width=w, height=h)
    if img_ar > box_ar:      # imagen más ancha → recorta laterales
        frac = (1 - box_ar / img_ar) / 2
        pic.crop_left = pic.crop_right = frac
    else:                    # imagen más alta → recorta arriba/abajo
        frac = (1 - img_ar / box_ar) / 2
        pic.crop_top = pic.crop_bottom = frac
    return pic


def footer(sl, dark=False, mixed=False):
    """mixed=True: texto sobre fondo claro, logo sobre fondo oscuro (divisores)."""
    c = INK if (mixed or not dark) else WHITE
    txt(sl, Inches(0.42), Inches(5.16), Inches(4), Inches(0.25),
        [("Intothecom ", {"bold": True, "color": c}),
         ("|Agencia Marketing Digital", {"color": ORANGE})], size=9)
    logo = LOGO_W if (dark or mixed) else LOGO_B
    sl.shapes.add_picture(logo, Inches(9.12), Inches(4.98), height=Inches(0.42))


def navbar(sl, active=None, dark_from=None):
    """Barra de navegación superior. dark_from = índice desde el que el fondo es oscuro."""
    xs, gap = Inches(1.30), Inches(0.10)
    x = xs
    for i, item in enumerate(NAV):
        wdt = Inches(0.16 + 0.075 * len(item))
        on_dark = dark_from is not None and i >= dark_from
        col = WHITE if on_dark else INK
        t = txt(sl, x, Inches(0.30), wdt, Inches(0.24), item, size=8.5,
                color=ORANGE if item == active else col, align=PP_ALIGN.CENTER)
        if item == active:
            ln = rect(sl, x + Inches(0.02), Inches(0.545), wdt - Inches(0.04),
                      Pt(1.6), fill=ORANGE)
        x += wdt + gap


_FONT_CACHE = {}
_FONT_PATHS = {True: "/root/.fonts/BeVietnamPro-Bold.ttf",
               False: "/root/.fonts/BeVietnamPro-Regular.ttf"}


def _measure(text, size_pt, bold):
    """Ancho real del texto en pulgadas, medido con la fuente instalada."""
    from PIL import ImageFont
    key = (size_pt, bold)
    if key not in _FONT_CACHE:
        # 1 pt = 1/72 in; renderizamos a 600 ppp para precisión
        px = max(1, int(round(size_pt / 72.0 * 600)))
        _FONT_CACHE[key] = ImageFont.truetype(_FONT_PATHS[bold], px)
    return _FONT_CACHE[key].getlength(text) / 600.0


def est_lines(text, width_in, size_pt, bold=True):
    """Número de líneas que ocupará el texto dentro de un ancho dado."""
    if not text:
        return 0
    try:
        total = 0
        for chunk in str(text).split("\n"):
            words, line, lines = chunk.split(), "", 1
            for wd in words:
                cand = wd if not line else line + " " + wd
                if _measure(cand, size_pt, bold) <= width_in:
                    line = cand
                else:
                    lines += 1
                    line = wd
            total += lines
        return total
    except Exception:
        per_char = size_pt * 0.62 / 72.0
        cap = max(1, int(width_in / per_char))
        return max(1, -(-len(str(text)) // cap))


def title2(sl, x, y, l1, l2, size=27, w=Inches(5.4), color1=INK, color2=ORANGE, gap=0.44):
    txt(sl, x, y, w, Inches(0.5), l1, size=size, bold=True, color=color1)
    txt(sl, x, y + Inches(gap), w, Inches(0.5), l2, size=size, bold=True, color=color2)


def title1(sl, x, y, text, size=27, w=Inches(8.6), color=INK, align=PP_ALIGN.LEFT):
    return txt(sl, x, y, w, Inches(0.6), text, size=size, bold=True, color=color, align=align)


def accent_arrow(sl, x, y, h=Inches(0.78)):
    """Línea vertical naranja con punta, como en el deck base."""
    rect(sl, x, y, Pt(1.4), h, fill=ORANGE)


def pill(sl, x, y, w, h, label, fill=ORANGE, color=WHITE, size=10.5, bold=True,
         radius=True):
    s = rect(sl, x, y, w, h, fill=fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE)
    try:
        s.adjustments[0] = 0.5 if radius else 0
    except Exception:
        pass
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return s


def card(sl, x, y, w, h, fill=WHITE, radius=0.05):
    s = rect(sl, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s
