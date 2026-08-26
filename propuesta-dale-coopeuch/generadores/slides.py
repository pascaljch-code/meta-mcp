# -*- coding: utf-8 -*-
"""Tipos de lámina que replican los layouts del deck base de Intothecom."""
from deck_lib import *
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def s_cover(prs, cliente, servicios):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=BLACK)
    picture_cover(sl, f"{PHOTO}/cover_laptop.jpeg", Inches(-0.35), Inches(0.35),
                  Inches(5.7), Inches(4.9))
    x = Inches(5.05)
    txt(sl, x, Inches(0.78), Inches(4.7), Inches(0.9), "PROPUESTA",
        size=45, bold=True, color=ORANGE)
    txt(sl, x, Inches(1.62), Inches(4.7), Inches(0.7), "DE SERVICIOS",
        size=31, bold=True, color=WHITE)
    txt(sl, x, Inches(2.42), Inches(4.7), Inches(0.6), cliente,
        size=27, bold=False, color=WHITE)
    txt(sl, x, Inches(3.30), Inches(4.6), Inches(0.6), servicios, size=9.5, color=LINE)
    pill(sl, x, Inches(3.95), Inches(1.62), Inches(0.36), "Comenzar  →", size=9.5)
    footer(sl, dark=True)
    return sl


def s_divider(prs, kicker, bajada, photo, active=None):
    """Divisor de sección: foto a la izquierda, panel oscuro a la derecha."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    rect(sl, Inches(4.33), 0, Inches(5.67), H, fill=INK)
    picture_cover(sl, photo, Inches(-0.35), Inches(1.02), Inches(5.35), Inches(3.62))
    navbar(sl, active=active, dark_from=3)
    x = Inches(5.60)
    txt(sl, x, Inches(1.72), Inches(4.2), Inches(1.0), kicker,
        size=25, bold=True, color=ORANGE, spacing=1.05)
    txt(sl, x, Inches(2.78), Inches(3.95), Inches(1.2), bajada, size=9.5, color=LINE,
        spacing=1.35)
    pill(sl, x, Inches(3.98), Inches(1.55), Inches(0.36), "Continuar  →", size=9.5)
    footer(sl, mixed=True)
    return sl


def s_split(prs, active, l1, l2, body=None, items=None, photo=None,
            photo_side="right", accent=True, item_size=10.5):
    """Título a dos tonos + texto/bullets a un lado y foto al otro."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    if photo:
        if photo_side == "right":
            picture_cover(sl, photo, Inches(5.19), Inches(0.72), Inches(4.81), Inches(3.63))
            rect(sl, Inches(5.19), Inches(4.35), Inches(4.81), Inches(1.275), fill=ORANGE)
            tx, tw = Inches(0.72), Inches(4.2)
        else:
            picture_cover(sl, photo, 0, Inches(0.72), Inches(4.55), Inches(3.63))
            rect(sl, 0, Inches(4.35), Inches(4.55), Inches(0.67), fill=ORANGE)
            tx, tw = Inches(5.05), Inches(4.4)
    else:
        tx, tw = Inches(0.72), Inches(8.6)
    navbar(sl, active=active)
    if accent:
        accent_arrow(sl, tx - Inches(0.26), Inches(0.95))
    tw_in = tw / 914400
    n_lines = est_lines(l1, tw_in, 23) + est_lines(l2, tw_in, 23)
    title2(sl, tx, Inches(0.92), l1, l2, size=23, w=tw,
           gap=0.44 * est_lines(l1, tw_in, 23))
    y = Inches(0.92 + 0.44 * n_lines + 0.20)
    if body:
        nb = est_lines(body, tw_in, 10, bold=False)
        blanks = sum(1 for ln in str(body).split("\n") if not ln.strip())
        txt(sl, tx, y, tw, Inches(0.21 * nb), body, size=10, color=GRAY, spacing=1.3)
        y += Inches(0.205 * nb + 0.06 * blanks + 0.24)
    if items:
        bullets(sl, tx, y, tw, Inches(2.4), items, size=item_size)
    footer(sl, mixed=bool(photo and photo_side == "right"))
    return sl


def s_numbered(prs, active, title_l1, title_l2, entries, intro=None, cols=1,
               photos=None):
    """Grilla con números naranjos grandes 01. 02. 03."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    navbar(sl, active=active)
    title2(sl, Inches(0.72), Inches(0.92), title_l1, title_l2, size=23, w=Inches(3.6))
    if intro:
        txt(sl, Inches(0.72), Inches(2.05), Inches(3.4), Inches(1.2), intro,
            size=9.5, color=GRAY, spacing=1.3)
    x0, y0 = Inches(4.55), Inches(0.95)
    step = Inches(1.36) if len(entries) <= 3 else Inches(1.02)
    for i, (num, head, desc) in enumerate(entries):
        y = y0 + step * i
        txt(sl, x0, y, Inches(0.78), Inches(0.4), num, size=19, bold=True, color=ORANGE)
        txt(sl, x0 + Inches(0.85), y + Inches(0.02), Inches(4.3), Inches(0.3), head,
            size=12, bold=True, color=ORANGE)
        txt(sl, x0 + Inches(0.85), y + Inches(0.30), Inches(4.35), Inches(0.85), desc,
            size=8.8, color=GRAY, spacing=1.25)
    if photos:
        if len(photos) == 1:
            picture_cover(sl, photos[0], Inches(0.72), Inches(2.72), Inches(3.5),
                          Inches(2.05))
        else:
            picture_cover(sl, photos[0], Inches(0.72), Inches(3.10), Inches(1.62),
                          Inches(1.68))
            picture_cover(sl, photos[1], Inches(2.46), Inches(2.62), Inches(1.76),
                          Inches(2.16))
    footer(sl)
    return sl


def s_grid(prs, active, title, cards, dark_panel=True, cols=3, sub=None):
    """Panel oscuro con celdas numeradas (layout 'Funciones Claves')."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    navbar(sl, active=active)
    title1(sl, Inches(0.6), Inches(0.88), title, size=23, w=Inches(8.8),
           align=PP_ALIGN.CENTER)
    if sub:
        txt(sl, Inches(1.2), Inches(1.42), Inches(7.6), Inches(0.3), sub, size=9.5,
            color=GRAY, align=PP_ALIGN.CENTER)
    top = Inches(1.78) if not sub else Inches(1.86)
    panel_h = Inches(3.18) if not sub else Inches(3.10)
    if dark_panel:
        card(sl, Inches(0.1), top, Inches(9.8), panel_h, fill=INK, radius=0.035)
    rows = (len(cards) + cols - 1) // cols
    cw = Inches(9.0 / cols)
    ch = (panel_h - Inches(0.4)) / rows
    for i, (head, desc) in enumerate(cards):
        r, c = divmod(i, cols)
        x = Inches(0.5) + cw * c
        y = top + Inches(0.24) + ch * r
        hw = (cw - Inches(0.72)) / 914400
        nh = est_lines(head, hw, 11.5)
        txt(sl, x, y, cw - Inches(0.72), Inches(0.3 * nh), head, size=11.5, bold=True,
            color=WHITE, spacing=1.1)
        txt(sl, x, y + Inches(0.20 + 0.22 * nh), cw - Inches(0.55), Inches(0.9), desc,
            size=8.6, color=LINE, spacing=1.25)
        pill(sl, x + cw - Inches(0.66), y - Inches(0.04), Inches(0.34), Inches(0.30),
             f"{i+1:02d}", size=9, radius=False)
    footer(sl)
    return sl


def s_stats(prs, active, title_l1, title_l2, body, stats, photo=None, brand=None):
    """Tarjetas de métricas; la primera en naranjo (layout 'Resultados Globales')."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    navbar(sl, active=active)
    txt(sl, Inches(0.6), Inches(0.88), Inches(8.8), Inches(0.5),
        [(title_l1 + " ", {"color": INK}), (title_l2, {"color": ORANGE})],
        size=23, bold=True)
    txt(sl, Inches(0.6), Inches(1.52), Inches(4.6), Inches(1.1), body, size=9,
        color=GRAY, spacing=1.3)
    if photo:
        picture_cover(sl, photo, Inches(0.6), Inches(2.55), Inches(4.35), Inches(2.35))
    if brand:
        tb = txt(sl, Inches(-0.62), Inches(3.35), Inches(1.9), Inches(0.3), brand, size=9,
                 color=INK, align=PP_ALIGN.CENTER)
        tb.rotation = 270
    x0, y0 = Inches(5.35), Inches(1.45)
    cw, chh, g = Inches(2.05), Inches(1.55), Inches(0.18)
    for i, (big, small) in enumerate(stats[:4]):
        r, c = divmod(i, 2)
        x = x0 + (cw + g) * c
        y = y0 + (chh + g) * r
        is_first = (i == 0)
        card(sl, x, y, cw, chh, fill=ORANGE if is_first else WHITE, radius=0.08)
        if not is_first:
            card(sl, x, y, cw, chh, fill=RGBColor(0xFA, 0xFA, 0xFA), radius=0.08)
        txt(sl, x, y + Inches(0.34), cw, Inches(0.5), big, size=25, bold=True,
            color=WHITE if is_first else INK, align=PP_ALIGN.CENTER)
        txt(sl, x + Inches(0.1), y + Inches(0.92), cw - Inches(0.2), Inches(0.5), small,
            size=9.5, color=WHITE if is_first else GRAY_D, align=PP_ALIGN.CENTER,
            spacing=1.15)
    footer(sl)
    return sl


def s_spec(prs, active, l1, l2, items, photo, tag):
    """Lámina de especificación de servicio: bullets + foto vertical + tag lateral."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    rect(sl, Inches(6.55), Inches(3.05), Inches(3.45), Inches(2.575), fill=ORANGE)
    picture_cover(sl, photo, Inches(5.15), Inches(0.72), Inches(4.0), Inches(4.13))
    navbar(sl, active=active)
    n1 = est_lines(l1, 4.4, 22)
    n_lines = n1 + est_lines(l2, 4.4, 22)
    title2(sl, Inches(0.58), Inches(0.90), l1, l2, size=22, w=Inches(4.4),
           color1=ORANGE, color2=INK, gap=0.42 * n1)
    bullets(sl, Inches(0.58), Inches(0.90 + 0.42 * n_lines + 0.20), Inches(4.35),
            Inches(3.0), items, size=9.5, gap=6)
    tb = txt(sl, Inches(9.15), Inches(3.15), Inches(0.7), Inches(1.9), tag, size=13,
             color=INK, bold=False)
    tb.rotation = 270
    footer(sl)
    return sl


def s_price(prs, cards_data, subtitle="Servicios Mensuales", note=None):
    """Propuesta comercial: tarjetas blancas sobre banda oscura."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    rect(sl, 0, Inches(2.72), W, Inches(2.905), fill=INK)
    navbar(sl, active="Propuesta Com.")
    title1(sl, Inches(0.6), Inches(0.80), "Propuesta Comercial", size=26,
           w=Inches(8.8), align=PP_ALIGN.CENTER)
    txt(sl, Inches(0.6), Inches(1.28), Inches(8.8), Inches(0.3), subtitle, size=13,
        color=INK, align=PP_ALIGN.CENTER)
    n = len(cards_data)
    total_w = Inches(9.2)
    cw = Inches((9.2 - 0.22 * (n - 1)) / n)
    x = Inches(0.4)
    for name, feats, valor, iva, total in cards_data:
        card(sl, x, Inches(1.72), cw, Inches(3.30), fill=WHITE, radius=0.045)
        pill(sl, x + Inches(0.14), Inches(1.90), cw - Inches(0.28), Inches(0.33), name,
             size=10)
        bullets(sl, x + Inches(0.20), Inches(2.42), cw - Inches(0.38), Inches(1.5),
                feats, size=8.2, color=GRAY_D, bullet_color=GRAY, gap=3, spacing=1.15)
        yv = Inches(4.02)
        rect(sl, x + Inches(0.20), yv - Inches(0.06), cw - Inches(0.4), Pt(1.1),
             fill=ORANGE)
        for i, (k, v) in enumerate([("VALOR", valor), ("IVA", iva)]):
            txt(sl, x + Inches(0.22), yv + Inches(0.06 + 0.21 * i), Inches(1.0),
                Inches(0.2), k, size=9, bold=True, color=GRAY_D)
            txt(sl, x + Inches(0.20), yv + Inches(0.06 + 0.21 * i), cw - Inches(0.42),
                Inches(0.2), v, size=9, bold=True, color=INK, align=PP_ALIGN.RIGHT)
        yt = yv + Inches(0.50)
        rect(sl, x + Inches(0.20), yt, cw - Inches(0.4), Pt(1.1), fill=ORANGE)
        txt(sl, x + Inches(0.22), yt + Inches(0.09), Inches(1.0), Inches(0.2), "TOTAL",
            size=9.5, bold=True, color=INK)
        txt(sl, x + Inches(0.20), yt + Inches(0.09), cw - Inches(0.42), Inches(0.2),
            total, size=9.5, bold=True, color=INK, align=PP_ALIGN.RIGHT)
        x += cw + Inches(0.22)
    if note:
        txt(sl, Inches(3.15), Inches(5.17), Inches(5.85), Inches(0.25), note, size=6.8,
            color=RGBColor(0xB0, 0xB0, 0xB0), italic=True)
    footer(sl, dark=True)
    return sl


def s_text(prs, active, title, items, sub=None, two_col=False, size=10):
    """Lámina de texto simple, una o dos columnas."""
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    navbar(sl, active=active)
    accent_arrow(sl, Inches(0.46), Inches(0.95))
    title1(sl, Inches(0.72), Inches(0.90), title, size=23, w=Inches(8.6))
    y = Inches(1.62)
    if sub:
        txt(sl, Inches(0.72), y, Inches(8.5), Inches(0.5), sub, size=9.5, color=GRAY,
            spacing=1.3)
        y += Inches(0.52)
    if two_col:
        half = (len(items) + 1) // 2
        bullets(sl, Inches(0.72), y, Inches(4.1), Inches(3.0), items[:half], size=size,
                gap=6)
        bullets(sl, Inches(5.15), y, Inches(4.1), Inches(3.0), items[half:], size=size,
                gap=6)
    else:
        bullets(sl, Inches(0.72), y, Inches(8.5), Inches(3.2), items, size=size, gap=6)
    footer(sl)
    return sl


def s_contact(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    rect(sl, 0, Inches(3.95), W, Inches(1.675), fill=INK)
    navbar(sl)
    picture_cover(sl, f"{PHOTO}/office_wide.jpeg", Inches(5.2), Inches(0.85),
                  Inches(4.3), Inches(3.9))
    sl.shapes.add_picture(LOGO_W, Inches(6.35), Inches(1.85), height=Inches(1.35))
    txt(sl, Inches(0.65), Inches(1.55), Inches(4.2), Inches(0.7), "Contacto", size=34,
        bold=True, color=ORANGE)
    txt(sl, Inches(0.65), Inches(2.25), Inches(3.9), Inches(1.1),
        [("Estamos aquí para ayudarte. Si tienes dudas, consultas o necesitas más "
          "información sobre nuestros servicios, contáctanos. Nuestro equipo está listo "
          "para ", {}),
         ("asesorarte y potenciar tu marca con una estrategia digital a medida.",
          {"bold": True, "color": INK})], size=9.5, color=GRAY, spacing=1.35)
    b = card(sl, Inches(0.3), Inches(3.45), Inches(5.2), Inches(1.05), fill=ORANGE,
             radius=0.08)
    txt(sl, Inches(0.62), Inches(3.63), Inches(2.4), Inches(0.25), "Email & Web",
        size=11.5, bold=True, color=WHITE)
    txt(sl, Inches(0.62), Inches(3.90), Inches(2.6), Inches(0.5),
        "ignacio@intothecom.com\nwww.intothecom.com", size=8.5, color=WHITE,
        spacing=1.25)
    txt(sl, Inches(3.35), Inches(3.63), Inches(2.0), Inches(0.25), "Teléfono",
        size=11.5, bold=True, color=WHITE)
    txt(sl, Inches(3.35), Inches(3.90), Inches(2.0), Inches(0.25), "+56 9 5016 0966",
        size=8.5, color=WHITE)
    footer(sl, dark=True)
    return sl
